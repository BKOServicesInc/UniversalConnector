"""
UniversalConnector — Architecture Diagram Generator
Produces a polished .pptx with a single landscape slide.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from lxml import etree
import copy

# ── Palette ──────────────────────────────────────────────────────────────────
DARK_BG      = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
PANEL_BG     = RGBColor(0x10, 0x28, 0x3C)   # dark teal-navy
CARD_BLUE    = RGBColor(0x1A, 0x4A, 0x7A)   # medium blue  (adapters)
CARD_TEAL    = RGBColor(0x0D, 0x6E, 0x6E)   # teal         (engine)
CARD_PURPLE  = RGBColor(0x4A, 0x2A, 0x7A)   # purple       (nats/sink)
CARD_GREEN   = RGBColor(0x14, 0x53, 0x38)   # dark green   (sources)
CARD_ORANGE  = RGBColor(0x7A, 0x3C, 0x0A)   # amber        (outputs)
ACCENT       = RGBColor(0x00, 0xD4, 0xFF)   # cyan accent
ACCENT2      = RGBColor(0xFF, 0xA5, 0x00)   # orange accent
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY   = RGBColor(0xBB, 0xCC, 0xDD)
MID_GREY     = RGBColor(0x55, 0x77, 0x99)

# ── Slide dimensions (16:9 widescreen) ───────────────────────────────────────
SW = Inches(13.33)
SH = Inches(7.5)


def rgb_hex(c: RGBColor) -> str:
    return f"{c[0]:02X}{c[1]:02X}{c[2]:02X}"


def add_rect(slide, x, y, w, h, fill: RGBColor, alpha=None,
             line_color: RGBColor = None, line_width_pt=0.75, radius_emu=0):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE — rounded via XML tweak below
        x, y, w, h
    )
    fill_fmt = shape.fill
    fill_fmt.solid()
    fill_fmt.fore_color.rgb = fill

    line = shape.line
    if line_color:
        line.color.rgb = line_color
        line.width = Pt(line_width_pt)
    else:
        line.fill.background()  # no border

    if radius_emu:
        # Patch prstGeom to use roundRect
        sp = shape._element
        spPr = sp.find(qn('p:spPr'))
        prstGeom = spPr.find(qn('a:prstGeom'))
        if prstGeom is not None:
            prstGeom.set('prst', 'roundRect')
            avLst = prstGeom.find(qn('a:avLst'))
            if avLst is None:
                avLst = etree.SubElement(prstGeom, qn('a:avLst'))
            else:
                avLst.clear()
            gd = etree.SubElement(avLst, qn('a:gd'))
            # adj value: 0–50000 (50000 = fully round). ~8000 gives a nice small radius.
            gd.set('name', 'adj')
            gd.set('fmla', f'val {radius_emu}')
    return shape


def add_label(slide, text, x, y, w, h,
              font_size=10, bold=False, color: RGBColor = WHITE,
              align=PP_ALIGN.CENTER, italic=False):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_connector_arrow(slide, x1, y1, x2, y2,
                        color: RGBColor = ACCENT, width_pt=1.5, dashed=False):
    """Draw a line with an arrowhead at the end."""
    cx = slide.shapes._spTree
    # Build a cxnSp (connector shape) via direct XML
    nvCxnSpPr = f'''<p:cxnSp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
              xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
              xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvCxnSpPr>
    <p:cNvPr id="9999" name="Arrow"/>
    <p:cNvCxnSpPr/>
    <p:nvPr/>
  </p:nvCxnSpPr>
  <p:spPr>
    <a:xfrm>
      <a:off x="{min(x1,x2)}" y="{min(y1,y2)}"/>
      <a:ext cx="{abs(x2-x1)}" cy="{abs(y2-y1)}"/>
    </a:xfrm>
    <a:prstGeom prst="line"><a:avLst/></a:prstGeom>
    <a:ln w="{int(width_pt * 12700)}">
      <a:solidFill><a:srgbClr val="{rgb_hex(color)}"/></a:solidFill>
      {'<a:prstDash val="dash"/>' if dashed else ''}
      <a:tailEnd type="arrow" w="med" len="med"/>
    </a:ln>
  </p:spPr>
</p:cxnSp>'''
    cx.append(etree.fromstring(cx_ns_fix(cx_ns_fix(nvCxnSpPr))))


def cx_ns_fix(xml_str):
    """Fix namespace prefix for drawingml — python-pptx uses 'a:' for drawingml."""
    return xml_str.replace(
        'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"',
        'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
    )


# ─────────────────────────────────────────────────────────────────────────────
def build_slide(prs):
    slide_layout = prs.slide_layouts[6]   # blank
    slide = prs.slides.add_slide(slide_layout)

    # ── Background ───────────────────────────────────────────────────────────
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG

    # ── Title banner ─────────────────────────────────────────────────────────
    banner = add_rect(slide,
                      x=Inches(0.25), y=Inches(0.15),
                      w=Inches(12.83), h=Inches(0.55),
                      fill=PANEL_BG, line_color=ACCENT, line_width_pt=1.5,
                      radius_emu=6000)
    add_label(slide, "UniversalConnector  —  Architecture Overview",
              x=Inches(0.25), y=Inches(0.18),
              w=Inches(12.83), h=Inches(0.50),
              font_size=18, bold=True, color=ACCENT)

    add_label(slide, ".NET 10 Worker Service  ·  NATS JetStream  ·  PostgreSQL CDC  ·  Descriptor-Driven",
              x=Inches(0.25), y=Inches(0.62),
              w=Inches(12.83), h=Inches(0.22),
              font_size=9, bold=False, color=MID_GREY, italic=True)

    # ═══════════════════════════════════════════════════════════════════════
    # COLUMN 1 — Data Sources   (x ≈ 0.25)
    # ═══════════════════════════════════════════════════════════════════════
    col1_x = Inches(0.25)
    col1_w = Inches(2.10)
    panel1 = add_rect(slide,
                      x=col1_x, y=Inches(0.9),
                      w=col1_w, h=Inches(6.1),
                      fill=PANEL_BG, line_color=MID_GREY, line_width_pt=0.75,
                      radius_emu=5000)
    add_label(slide, "DATA  SOURCES",
              x=col1_x, y=Inches(0.92),
              w=col1_w, h=Inches(0.3),
              font_size=9, bold=True, color=ACCENT2)

    sources = [
        ("PostgreSQL",    "CDC  (WAL)  /  Polling"),
        ("SQL Server",    "Change Tracking  /  Polling"),
        ("MongoDB",       "Change Streams  /  Polling"),
        ("Neo4j",         "Streaming  /  Polling"),
        ("Databricks",    "Delta  /  Polling"),
        ("HTTP / REST",   "Polling  /  Webhooks"),
        ("Custom",        "IProtocolAdapter"),
    ]
    src_h = Inches(0.62)
    for i, (name, sub) in enumerate(sources):
        sy = Inches(1.30) + i * (src_h + Inches(0.08))
        add_rect(slide,
                 x=col1_x + Inches(0.1), y=sy,
                 w=col1_w - Inches(0.2), h=src_h,
                 fill=CARD_GREEN, line_color=ACCENT, line_width_pt=0.5,
                 radius_emu=4000)
        add_label(slide, name,
                  x=col1_x + Inches(0.12), y=sy + Inches(0.05),
                  w=col1_w - Inches(0.24), h=Inches(0.28),
                  font_size=10, bold=True, color=WHITE)
        add_label(slide, sub,
                  x=col1_x + Inches(0.12), y=sy + Inches(0.30),
                  w=col1_w - Inches(0.24), h=Inches(0.25),
                  font_size=7.5, bold=False, color=LIGHT_GREY)

    # ═══════════════════════════════════════════════════════════════════════
    # COLUMN 2 — Adapter Layer  (x ≈ 2.60)
    # ═══════════════════════════════════════════════════════════════════════
    col2_x = Inches(2.60)
    col2_w = Inches(2.20)
    panel2 = add_rect(slide,
                      x=col2_x, y=Inches(0.9),
                      w=col2_w, h=Inches(6.1),
                      fill=PANEL_BG, line_color=MID_GREY, line_width_pt=0.75,
                      radius_emu=5000)
    add_label(slide, "ADAPTER  LAYER",
              x=col2_x, y=Inches(0.92),
              w=col2_w, h=Inches(0.3),
              font_size=9, bold=True, color=ACCENT2)

    add_label(slide, "IProtocolAdapter",
              x=col2_x, y=Inches(1.22),
              w=col2_w, h=Inches(0.22),
              font_size=8, bold=False, color=MID_GREY, italic=True)

    adapters = [
        "PostgresAdapter",
        "SqlServerAdapter",
        "MongoDbAdapter",
        "Neo4jAdapter",
        "DatabricksAdapter",
        "HttpAdapter",
        "CustomAdapter",
    ]
    adp_h = Inches(0.62)
    for i, name in enumerate(adapters):
        ay = Inches(1.48) + i * (adp_h + Inches(0.08))
        add_rect(slide,
                 x=col2_x + Inches(0.1), y=ay,
                 w=col2_w - Inches(0.2), h=adp_h,
                 fill=CARD_BLUE, line_color=ACCENT, line_width_pt=0.5,
                 radius_emu=4000)
        add_label(slide, name,
                  x=col2_x + Inches(0.12), y=ay + Inches(0.12),
                  w=col2_w - Inches(0.24), h=Inches(0.38),
                  font_size=9.5, bold=True, color=WHITE)

    # ═══════════════════════════════════════════════════════════════════════
    # COLUMN 3 — Generic Engine  (x ≈ 5.05)
    # ═══════════════════════════════════════════════════════════════════════
    col3_x = Inches(5.05)
    col3_w = Inches(2.70)
    panel3 = add_rect(slide,
                      x=col3_x, y=Inches(0.9),
                      w=col3_w, h=Inches(6.1),
                      fill=PANEL_BG, line_color=MID_GREY, line_width_pt=0.75,
                      radius_emu=5000)
    add_label(slide, "GENERIC  ENGINE",
              x=col3_x, y=Inches(0.92),
              w=col3_w, h=Inches(0.3),
              font_size=9, bold=True, color=ACCENT2)

    engine_items = [
        ("DescriptorLoader", "YAML / JSON → ConnectorDescriptor"),
        ("DescriptorValidator", "Schema + constraint checks"),
        ("DescriptorBootstrapService", "IHostedService · startup"),
        ("ConnectorRegistry", "connectorId → IConnector"),
        ("GenericConnector", "BaseConnector + retry loop"),
        ("Snapshot Cache", "previous_payload  (polling)"),
        ("AdapterRegistry", "sourceType → IProtocolAdapter"),
        ("MultiSourceGenericFactory", "Descriptor → Connector factory"),
    ]
    eng_h = Inches(0.58)
    for i, (name, sub) in enumerate(engine_items):
        ey = Inches(1.28) + i * (eng_h + Inches(0.07))
        add_rect(slide,
                 x=col3_x + Inches(0.1), y=ey,
                 w=col3_w - Inches(0.2), h=eng_h,
                 fill=CARD_TEAL, line_color=ACCENT, line_width_pt=0.5,
                 radius_emu=4000)
        add_label(slide, name,
                  x=col3_x + Inches(0.14), y=ey + Inches(0.04),
                  w=col3_w - Inches(0.28), h=Inches(0.26),
                  font_size=9, bold=True, color=WHITE)
        add_label(slide, sub,
                  x=col3_x + Inches(0.14), y=ey + Inches(0.28),
                  w=col3_w - Inches(0.28), h=Inches(0.22),
                  font_size=7, bold=False, color=LIGHT_GREY)

    # ═══════════════════════════════════════════════════════════════════════
    # COLUMN 4 — Host Layer  (x ≈ 8.00)
    # ═══════════════════════════════════════════════════════════════════════
    col4_x = Inches(8.00)
    col4_w = Inches(2.35)
    panel4 = add_rect(slide,
                      x=col4_x, y=Inches(0.9),
                      w=col4_w, h=Inches(6.1),
                      fill=PANEL_BG, line_color=MID_GREY, line_width_pt=0.75,
                      radius_emu=5000)
    add_label(slide, "HOST  LAYER",
              x=col4_x, y=Inches(0.92),
              w=col4_w, h=Inches(0.3),
              font_size=9, bold=True, color=ACCENT2)

    host_items = [
        ("ConnectorPipelineService", "IHostedService\nOrchestrates all connectors"),
        ("NatsPublisher", "INatsPublisher\nNATS.Net 2.x publish"),
        ("PostgresDataSink", "IDataSink\nDapper INSERT"),
        ("ServiceCollectionExtensions", "DI wiring\nAll registrations"),
        ("appsettings.json", "NatsOptions\nPostgresSinkOptions"),
    ]
    host_h = Inches(0.80)
    for i, (name, sub) in enumerate(host_items):
        hy = Inches(1.28) + i * (host_h + Inches(0.12))
        add_rect(slide,
                 x=col4_x + Inches(0.1), y=hy,
                 w=col4_w - Inches(0.2), h=host_h,
                 fill=CARD_PURPLE, line_color=ACCENT, line_width_pt=0.5,
                 radius_emu=4000)
        add_label(slide, name,
                  x=col4_x + Inches(0.14), y=hy + Inches(0.05),
                  w=col4_w - Inches(0.28), h=Inches(0.28),
                  font_size=9, bold=True, color=WHITE)
        add_label(slide, sub,
                  x=col4_x + Inches(0.14), y=hy + Inches(0.30),
                  w=col4_w - Inches(0.28), h=Inches(0.38),
                  font_size=7.5, bold=False, color=LIGHT_GREY)

    # ═══════════════════════════════════════════════════════════════════════
    # COLUMN 5 — Outputs  (x ≈ 10.60)
    # ═══════════════════════════════════════════════════════════════════════
    col5_x = Inches(10.60)
    col5_w = Inches(2.48)
    panel5 = add_rect(slide,
                      x=col5_x, y=Inches(0.9),
                      w=col5_w, h=Inches(6.1),
                      fill=PANEL_BG, line_color=MID_GREY, line_width_pt=0.75,
                      radius_emu=5000)
    add_label(slide, "OUTPUTS",
              x=col5_x, y=Inches(0.92),
              w=col5_w, h=Inches(0.3),
              font_size=9, bold=True, color=ACCENT2)

    # NATS block
    add_rect(slide,
             x=col5_x + Inches(0.1), y=Inches(1.25),
             w=col5_w - Inches(0.2), h=Inches(1.80),
             fill=CARD_ORANGE, line_color=ACCENT2, line_width_pt=1.0,
             radius_emu=4000)
    add_label(slide, "NATS JetStream",
              x=col5_x + Inches(0.12), y=Inches(1.28),
              w=col5_w - Inches(0.24), h=Inches(0.30),
              font_size=11, bold=True, color=ACCENT2)
    add_label(slide, "Subject pattern:",
              x=col5_x + Inches(0.12), y=Inches(1.58),
              w=col5_w - Inches(0.24), h=Inches(0.22),
              font_size=8, bold=False, color=LIGHT_GREY)
    add_label(slide, "{prefix}.{sourceType}",
              x=col5_x + Inches(0.12), y=Inches(1.76),
              w=col5_w - Inches(0.24), h=Inches(0.22),
              font_size=8, bold=True, color=WHITE)
    add_label(slide, ".{connectorId}.{changeType}",
              x=col5_x + Inches(0.12), y=Inches(1.94),
              w=col5_w - Inches(0.24), h=Inches(0.22),
              font_size=8, bold=True, color=WHITE)

    # Postgres sink block
    add_rect(slide,
             x=col5_x + Inches(0.1), y=Inches(3.20),
             w=col5_w - Inches(0.2), h=Inches(2.10),
             fill=CARD_ORANGE, line_color=ACCENT2, line_width_pt=1.0,
             radius_emu=4000)
    add_label(slide, "PostgreSQL  (CDCDB)",
              x=col5_x + Inches(0.12), y=Inches(3.23),
              w=col5_w - Inches(0.24), h=Inches(0.30),
              font_size=11, bold=True, color=ACCENT2)
    add_label(slide, "Table:  data_changes",
              x=col5_x + Inches(0.12), y=Inches(3.55),
              w=col5_w - Inches(0.24), h=Inches(0.22),
              font_size=8.5, bold=False, color=LIGHT_GREY)

    pg_cols = [
        "event_id  UUID  PK",
        "source_type  /  connector_id",
        "entity_path  /  change_type",
        "primary_key  JSONB  (GIN)",
        "payload  /  previous_payload",
        "metadata  /  sequence_number",
    ]
    for i, col in enumerate(pg_cols):
        add_label(slide, f"▸  {col}",
                  x=col5_x + Inches(0.14), y=Inches(3.78) + i * Inches(0.245),
                  w=col5_w - Inches(0.28), h=Inches(0.22),
                  font_size=7, bold=False, color=LIGHT_GREY)

    # ── DataChangeEvent label between host and outputs ────────────────────
    add_rect(slide,
             x=Inches(10.0), y=Inches(6.2),
             w=Inches(3.1), h=Inches(0.55),
             fill=CARD_TEAL, line_color=ACCENT, line_width_pt=0.8,
             radius_emu=5000)
    add_label(slide, "DataChangeEvent  (sealed record)",
              x=Inches(10.0), y=Inches(6.28),
              w=Inches(3.1), h=Inches(0.38),
              font_size=8.5, bold=True, color=WHITE)

    # ── Legend ───────────────────────────────────────────────────────────────
    legend_x = Inches(0.27)
    legend_y = Inches(7.08)
    items = [
        (CARD_GREEN,  "Data Source"),
        (CARD_BLUE,   "Protocol Adapter"),
        (CARD_TEAL,   "Engine Component"),
        (CARD_PURPLE, "Host / DI"),
        (CARD_ORANGE, "Output / Sink"),
    ]
    for i, (color, label) in enumerate(items):
        lx = legend_x + i * Inches(2.55)
        add_rect(slide, x=lx, y=legend_y, w=Inches(0.18), h=Inches(0.18),
                 fill=color, radius_emu=2000)
        add_label(slide, label,
                  x=lx + Inches(0.22), y=legend_y - Inches(0.01),
                  w=Inches(2.2), h=Inches(0.22),
                  font_size=8, bold=False, color=LIGHT_GREY,
                  align=PP_ALIGN.LEFT)

    # ── Flow arrows (simplified — horizontal between column midpoints) ────────
    # Sources → Adapters
    for i in range(7):
        y_mid = int(Inches(1.61) + i * int(Inches(0.70)))
        _draw_h_arrow(slide, int(Inches(2.37)), y_mid, int(Inches(2.60)), y_mid)

    # Adapters → Engine
    for i in range(7):
        y_mid = int(Inches(1.79) + i * int(Inches(0.70)))
        _draw_h_arrow(slide, int(Inches(4.82)), y_mid, int(Inches(5.05)), y_mid)

    # Engine → Host
    _draw_h_arrow(slide, int(Inches(7.77)), int(Inches(4.33)),
                  int(Inches(8.00)), int(Inches(4.33)))

    # Host → NATS
    _draw_h_arrow(slide, int(Inches(10.37)), int(Inches(2.15)),
                  int(Inches(10.60)), int(Inches(2.15)), color=ACCENT2)

    # Host → PG Sink
    _draw_h_arrow(slide, int(Inches(10.37)), int(Inches(4.25)),
                  int(Inches(10.60)), int(Inches(4.25)), color=ACCENT2)

    # ── Footer ───────────────────────────────────────────────────────────────
    add_label(slide,
              "UniversalConnector  ·  BKO Services  ·  .NET 10  ·  NATS.Net 2.x  ·  Npgsql 10.x  ·  2026",
              x=Inches(0.25), y=Inches(7.32),
              w=Inches(12.83), h=Inches(0.18),
              font_size=7, bold=False, color=MID_GREY, italic=True)

    return slide


def _draw_h_arrow(slide, x1, y1, x2, y2, color: RGBColor = ACCENT):
    """Draw a simple horizontal connector arrow."""
    cx = slide.shapes._spTree
    nsmap = {
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    }
    xml_str = f'''<p:cxnSp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
              xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
              xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvCxnSpPr>
    <p:cNvPr id="1" name="Connector"/>
    <p:cNvCxnSpPr/>
    <p:nvPr/>
  </p:nvCxnSpPr>
  <p:spPr>
    <a:xfrm>
      <a:off x="{min(x1,x2)}" y="{min(y1,y2)}"/>
      <a:ext cx="{abs(x2-x1) or 1}" cy="{abs(y2-y1) or 1}"/>
    </a:xfrm>
    <a:prstGeom prst="line"><a:avLst/></a:prstGeom>
    <a:ln w="19050">
      <a:solidFill><a:srgbClr val="{rgb_hex(color)}"/></a:solidFill>
      <a:tailEnd type="arrow" w="sm" len="sm"/>
    </a:ln>
  </p:spPr>
</p:cxnSp>'''
    cx.append(etree.fromstring(xml_str))


# ─────────────────────────────────────────────────────────────────────────────
def main():
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH

    build_slide(prs)

    out = r"C:\Repos\UniversalConnector\UniversalConnector_Architecture.pptx"
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == "__main__":
    main()
