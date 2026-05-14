"""
Generates UniversalConnector architecture diagram as a PowerPoint slide.
Run: py make_diagram.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

# ── Colour palette ──────────────────────────────────────────────────────────
DARK_BG      = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_TEAL  = RGBColor(0x00, 0xB4, 0xD8)
ACCENT_GREEN = RGBColor(0x06, 0xD6, 0xA0)
ACCENT_AMBER = RGBColor(0xFF, 0xB7, 0x03)
ACCENT_RED   = RGBColor(0xEF, 0x47, 0x6F)
ACCENT_PURP  = RGBColor(0x9B, 0x5D, 0xE5)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY   = RGBColor(0xCC, 0xD6, 0xE8)
YELLOW       = RGBColor(0xF7, 0xEC, 0x09)
LIME         = RGBColor(0x80, 0xFF, 0x80)

# RGBColor.__str__ returns 6-char hex e.g. "00B4D8"
def hex6(c): return str(c)

# ── Slide setup ──────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(16)
prs.slide_height = Inches(9)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = DARK_BG

# ── XML helpers ──────────────────────────────────────────────────────────────
NS = {
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
}

_id_counter = [100]

def next_id():
    _id_counter[0] += 1
    return _id_counter[0]

def _append(slide, xml_str):
    el = etree.fromstring(xml_str)
    slide.shapes._spTree.append(el)

def add_rect(slide, x, y, w, h,
             fill, border=None, border_pt=1.5, corner=15000):
    """Rounded rectangle. corner = adj value (0–50000)."""
    bid = next_id()
    bw  = int(Pt(border_pt)) if border else 0
    border_xml = (
        f'<a:ln w="{bw}" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<a:solidFill><a:srgbClr val="{hex6(border)}"/></a:solidFill></a:ln>'
        if border else
        '<a:ln w="0" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        '<a:noFill/></a:ln>'
    )
    xml = f'''<p:sp
        xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
        xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
      <p:nvSpPr>
        <p:cNvPr id="{bid}" name="r{bid}"/>
        <p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>
        <p:nvPr/>
      </p:nvSpPr>
      <p:spPr>
        <a:xfrm><a:off x="{int(Inches(x))}" y="{int(Inches(y))}"/>
                <a:ext cx="{int(Inches(w))}" cy="{int(Inches(h))}"/></a:xfrm>
        <a:prstGeom prst="roundRect">
          <a:avLst><a:gd name="adj" fmla="val {corner}"/></a:avLst>
        </a:prstGeom>
        <a:solidFill><a:srgbClr val="{hex6(fill)}"/></a:solidFill>
        {border_xml}
        <a:effectLst/>
      </p:spPr>
      <p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody>
    </p:sp>'''
    _append(slide, xml)


def add_text(slide, text, x, y, w, h,
             size=10, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_lines(slide, lines, x, y, w, h,
              size=7.5, color=LIGHT_GREY, bold_first=False):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        run.font.size  = Pt(size)
        run.font.bold  = bold_first and i == 0
        run.font.color.rgb = color


def add_arrow(slide, x1, y1, x2, y2, color=ACCENT_TEAL, pt=2.0):
    bid = next_id()
    cx  = int(abs(Inches(x2) - Inches(x1)))
    cy  = int(abs(Inches(y2) - Inches(y1)))
    ox  = int(min(Inches(x1), Inches(x2)))
    oy  = int(min(Inches(y1), Inches(y2)))
    # flip if going right-to-left or bottom-to-top
    flipH = 'flipH="1"' if x2 < x1 else ''
    flipV = 'flipV="1"' if y2 < y1 else ''
    xml = f'''<p:sp
        xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
        xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
      <p:nvSpPr>
        <p:cNvPr id="{bid}" name="arr{bid}"/>
        <p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>
        <p:nvPr/>
      </p:nvSpPr>
      <p:spPr>
        <a:xfrm {flipH} {flipV}>
          <a:off x="{ox}" y="{oy}"/>
          <a:ext cx="{max(cx,1)}" cy="{max(cy,1)}"/>
        </a:xfrm>
        <a:prstGeom prst="rightArrow">
          <a:avLst>
            <a:gd name="adj1" fmla="val 50000"/>
            <a:gd name="adj2" fmla="val 50000"/>
          </a:avLst>
        </a:prstGeom>
        <a:solidFill><a:srgbClr val="{hex6(color)}"/></a:solidFill>
        <a:ln w="0"><a:noFill/></a:ln>
      </p:spPr>
      <p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody>
    </p:sp>'''
    _append(slide, xml)


# ════════════════════════════════════════════════════════════════════════════
# LAYOUT CONSTANTS
# ════════════════════════════════════════════════════════════════════════════
TOP   = 0.80   # top of all column boxes
BOT   = 8.70   # bottom of all column boxes
H     = BOT - TOP

C1_X, C1_W = 0.15, 2.45   # Data Sources
C2_X, C2_W = 2.80, 3.20   # Generic Driver Layer
C3_X, C3_W = 6.20, 3.35   # Host / Orchestration
C4_X, C4_W = 9.75, 3.10   # Infrastructure
C5_X, C5_W = 13.05, 2.75  # NATS Broker

COL_BG_ALPHA = 0.12  # conceptual only — pptx has no alpha

# ── Title bar ────────────────────────────────────────────────────────────────
add_rect(slide, 0, 0, 16, 0.68, DARK_BG, border=None, corner=0)
add_rect(slide, 0, 0, 16, 0.68, RGBColor(0x0F, 0x3C, 0x78), corner=0)
add_text(slide, "UNIVERSAL CONNECTOR — ARCHITECTURE",
         0.25, 0.10, 11, 0.50, size=19, bold=True,
         color=ACCENT_TEAL, align=PP_ALIGN.LEFT)
add_text(slide, ".NET 10 Worker Service  •  CDC Bridge",
         10.5, 0.14, 5.3, 0.44, size=11, italic=True,
         color=LIGHT_GREY, align=PP_ALIGN.RIGHT)

# ════════════════════════════════════════════════════════════════════════════
# COLUMN 1 — DATA SOURCES
# ════════════════════════════════════════════════════════════════════════════
add_rect(slide, C1_X, TOP, C1_W, H,
         RGBColor(0x02, 0x2A, 0x1C), border=ACCENT_GREEN, border_pt=1.5, corner=8000)
add_rect(slide, C1_X, TOP, C1_W, 0.44,
         ACCENT_GREEN, corner=8000)
add_text(slide, "DATA SOURCES", C1_X+0.08, TOP+0.07, C1_W-0.15, 0.32,
         size=9, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)

# Databases group label
add_rect(slide, C1_X+0.12, TOP+0.55, C1_W-0.24, 0.28,
         RGBColor(0x04, 0x4A, 0x32), corner=5000)
add_text(slide, "Databases", C1_X+0.18, TOP+0.57, C1_W-0.35, 0.24,
         size=8, bold=True, color=ACCENT_GREEN)

db_items = [
    ("PostgreSQL",  "CDC (logical replication) / Polling"),
    ("SQL Server",  "CDC (change tracking) / Polling"),
    ("MongoDB",     "Change Streams / Polling"),
    ("Neo4j",       "Polling"),
    ("Databricks",  "CDC / Polling"),
]
for i, (name, mode) in enumerate(db_items):
    yy = TOP + 0.93 + i * 0.56
    add_rect(slide, C1_X+0.12, yy, C1_W-0.24, 0.50,
             RGBColor(0x02, 0x38, 0x27), border=ACCENT_GREEN, border_pt=0.6, corner=6000)
    add_text(slide, name, C1_X+0.20, yy+0.03, C1_W-0.40, 0.22,
             size=8, bold=True, color=WHITE)
    add_text(slide, mode, C1_X+0.20, yy+0.25, C1_W-0.40, 0.20,
             size=6.5, italic=True, color=ACCENT_GREEN)

# REST group label
add_rect(slide, C1_X+0.12, TOP+3.82, C1_W-0.24, 0.28,
         RGBColor(0x04, 0x4A, 0x32), corner=5000)
add_text(slide, "REST / HTTP APIs", C1_X+0.18, TOP+3.84, C1_W-0.35, 0.24,
         size=8, bold=True, color=ACCENT_GREEN)

api_items = [
    ("SharePoint",  "Delta / Polling"),
    ("SAP",         "Delta / Polling"),
    ("Seeq",        "Polling"),
    ("AVEVA PI",    "Polling"),
]
for i, (name, mode) in enumerate(api_items):
    yy = TOP + 4.20 + i * 0.54
    add_rect(slide, C1_X+0.12, yy, C1_W-0.24, 0.48,
             RGBColor(0x02, 0x38, 0x27), border=ACCENT_GREEN, border_pt=0.6, corner=6000)
    add_text(slide, name, C1_X+0.20, yy+0.03, C1_W-0.40, 0.22,
             size=8, bold=True, color=WHITE)
    add_text(slide, mode, C1_X+0.20, yy+0.25, C1_W-0.40, 0.20,
             size=6.5, italic=True, color=ACCENT_GREEN)

# Descriptor file hint
add_rect(slide, C1_X+0.12, TOP+6.52, C1_W-0.24, 0.56,
         RGBColor(0x01, 0x22, 0x18), border=ACCENT_GREEN, border_pt=1.0, corner=6000)
add_text(slide, "/connectors/", C1_X+0.18, TOP+6.56, C1_W-0.35, 0.22,
         size=8, bold=True, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)
add_text(slide, "*.yaml | *.json  descriptors", C1_X+0.18, TOP+6.76, C1_W-0.35, 0.22,
         size=7, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# COLUMN 2 — GENERIC DRIVER LAYER
# ════════════════════════════════════════════════════════════════════════════
add_rect(slide, C2_X, TOP, C2_W, H,
         RGBColor(0x06, 0x12, 0x2E), border=ACCENT_TEAL, border_pt=1.5, corner=8000)
add_rect(slide, C2_X, TOP, C2_W, 0.44,
         ACCENT_TEAL, corner=8000)
add_text(slide, "GENERIC DRIVER LAYER", C2_X+0.08, TOP+0.07, C2_W-0.15, 0.32,
         size=9, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)

def col2_block(y, h, title, lines=None, title_color=ACCENT_TEAL, border_color=None):
    bc = border_color or ACCENT_TEAL
    add_rect(slide, C2_X+0.15, y, C2_W-0.30, h,
             RGBColor(0x08, 0x1C, 0x48), border=bc, border_pt=0.8, corner=6000)
    add_text(slide, title, C2_X+0.24, y+0.05, C2_W-0.50, 0.26,
             size=8.5, bold=True, color=title_color)
    if lines:
        add_lines(slide, lines, C2_X+0.24, y+0.30, C2_W-0.50, h-0.32,
                  size=7.5, color=LIGHT_GREY)

col2_block(TOP+0.55, 0.58, "DescriptorBootstrapService",
           ["Loads all YAML/JSON files at startup",
            "Validates via DescriptorValidator  •  FailOnDescriptorError"],
           title_color=ACCENT_AMBER, border_color=ACCENT_AMBER)

col2_block(TOP+1.22, 0.58, "DescriptorValidator",
           ["sourceType ✓  mode ✓  connection fields ✓",
            "Warns: no entities, literal secrets, CDC slots"])

col2_block(TOP+1.90, 0.58, "ConnectorRegistry  +  MultiSourceGenericFactory",
           ["ResolveAll()  →  one GenericConnector per descriptor",
            "Factory lookup: sourceType → IProtocolAdapter"])

# GenericConnector — taller featured box
add_rect(slide, C2_X+0.15, TOP+2.58, C2_W-0.30, 1.40,
         RGBColor(0x04, 0x12, 0x38), border=ACCENT_TEAL, border_pt=1.5, corner=6000)
add_rect(slide, C2_X+0.15, TOP+2.58, C2_W-0.30, 0.32,
         ACCENT_TEAL, corner=6000)
add_text(slide, "GenericConnector  (extends BaseConnector)", C2_X+0.24, TOP+2.60, C2_W-0.50, 0.28,
         size=8.5, bold=True, color=DARK_BG)
add_lines(slide,
    ["ConnectAsync()      →  adapter.OpenAsync()",
     "StreamChangesAsync():",
     "   adapter yields RawChangeRecord",
     "   FieldMapper.Apply(record, descriptor.rules)",
     "   → yields RawChangeEvent  (normalized)",
     "DisconnectAsync()   →  adapter.CloseAsync()",
     "BaseConnector: failure tracking, 6-state machine"],
    C2_X+0.24, TOP+2.93, C2_W-0.50, 1.00, size=7.5)

col2_block(TOP+4.08, 0.52, "AdapterRegistry",
           ["Resolves IProtocolAdapter by sourceType at runtime"])

col2_block(TOP+4.70, 0.75, "FieldMapper",
           ["Rename  •  Type-cast  •  Exclude  •  Concept-map",
            "primaryKey dict  +  payload dict",
            "previous_fields  preserved on Update/Delete"])

col2_block(TOP+5.56, 0.52, "SubjectTemplateResolver",
           ["{context}.{entityPath}.{changeType}  →  NATS subject"])

# FSM box
add_rect(slide, C2_X+0.15, TOP+6.18, C2_W-0.30, 1.50,
         RGBColor(0x08, 0x1C, 0x48), border=ACCENT_TEAL, border_pt=0.8, corner=6000)
add_text(slide, "LifecycleFsm — Driver State Machine", C2_X+0.24, TOP+6.20, C2_W-0.50, 0.26,
         size=8.5, bold=True, color=ACCENT_TEAL)
add_lines(slide,
    ["Disconnected ─ start ──► Connecting",
     "Connecting   ─ connected ► Connected",
     "Connected    ─ streaming ► Streaming",
     "Streaming    ─ stop ───► Disconnected",
     "Streaming    ─ restart ─► Connecting",
     "Any state    ─ error ──► Failed"],
    C2_X+0.24, TOP+6.48, C2_W-0.50, 1.14, size=7.5)

# ════════════════════════════════════════════════════════════════════════════
# COLUMN 3 — HOST / ORCHESTRATION
# ════════════════════════════════════════════════════════════════════════════
add_rect(slide, C3_X, TOP, C3_W, H,
         RGBColor(0x12, 0x08, 0x2E), border=ACCENT_PURP, border_pt=1.5, corner=8000)
add_rect(slide, C3_X, TOP, C3_W, 0.44,
         ACCENT_PURP, corner=8000)
add_text(slide, "HOST / ORCHESTRATION", C3_X+0.08, TOP+0.07, C3_W-0.15, 0.32,
         size=9, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)

def col3_block(y, h, title, lines=None, title_color=ACCENT_PURP):
    add_rect(slide, C3_X+0.15, y, C3_W-0.30, h,
             RGBColor(0x1A, 0x0C, 0x3E), border=ACCENT_PURP, border_pt=0.8, corner=6000)
    add_text(slide, title, C3_X+0.24, y+0.05, C3_W-0.50, 0.26,
             size=8.5, bold=True, color=title_color)
    if lines:
        add_lines(slide, lines, C3_X+0.24, y+0.30, C3_W-0.50, h-0.32,
                  size=7.5, color=LIGHT_GREY)

col3_block(TOP+0.55, 0.58, "StartupSelfTestService",
           ["Creds file exists ✓  NATS ping ✓  JetStream ✓",
            "StopApplication() on critical failure"])

# ConnectorPipelineService — featured
add_rect(slide, C3_X+0.15, TOP+1.24, C3_W-0.30, 2.18,
         RGBColor(0x14, 0x06, 0x34), border=ACCENT_PURP, border_pt=1.5, corner=6000)
add_rect(slide, C3_X+0.15, TOP+1.24, C3_W-0.30, 0.32,
         ACCENT_PURP, corner=6000)
add_text(slide, "ConnectorPipelineService  (IHostedService)", C3_X+0.24, TOP+1.26, C3_W-0.50, 0.28,
         size=8.5, bold=True, color=DARK_BG)
add_lines(slide,
    ["1.  registry.ResolveAll()  →  list of ISourceDriver",
     "2.  LaunchDriverAsync(driver)  per driver:",
     "       ConnectAsync()",
     "       while not cancelled:",
     "           StreamChangesAsync()",
     "               pipeline.ProcessAsync(event)",
     "       DisconnectAsync()",
     "       on failure: exponential backoff restart",
     "       stops only if state == Failed",
     "3.  Start / Stop / Restart via lifecycle API",
     "4.  GetAllHealth()  →  per-driver HealthStatus"],
    C3_X+0.24, TOP+1.58, C3_W-0.50, 1.78, size=7.5)

# DefaultEventPipeline
add_rect(slide, C3_X+0.15, TOP+3.52, C3_W-0.30, 0.72,
         RGBColor(0x1A, 0x0C, 0x3E), border=ACCENT_PURP, border_pt=0.8, corner=6000)
add_text(slide, "DefaultEventPipeline  (IEventPipeline)", C3_X+0.24, TOP+3.54, C3_W-0.50, 0.26,
         size=8.5, bold=True, color=ACCENT_PURP)
add_lines(slide,
    ["1. NatsPublisher.PublishAsync(event)",
     "2. NatsCheckpointStore.SaveAsync(checkpoint)"],
    C3_X+0.24, TOP+3.80, C3_W-0.50, 0.38, size=7.5)

# Background services
add_rect(slide, C3_X+0.15, TOP+4.36, C3_W-0.30, 0.28,
         RGBColor(0x23, 0x10, 0x50), corner=4000)
add_text(slide, "Background Services  (always running)", C3_X+0.22, TOP+4.38, C3_W-0.40, 0.24,
         size=8, bold=True, color=ACCENT_PURP)

col3_block(TOP+4.74, 0.52, "HealthHeartbeatService",
           ["Publishes HealthStatus every 30 s  →  cdc.health.{driverId}"])
col3_block(TOP+5.36, 0.52, "DriverLifecycleService",
           ["Listens cdc.commands.*  •  dispatches start/stop/restart"])
col3_block(TOP+5.98, 0.52, "OntologyCacheRefreshService",
           ["Load Fuseki on startup  •  refresh on cdc.ontology.refresh"])
col3_block(TOP+6.60, 0.52, "NatsHealthCheck  (IHealthCheck)",
           ["3 s ping shared connection  →  Healthy / Unhealthy"])

# ════════════════════════════════════════════════════════════════════════════
# COLUMN 4 — INFRASTRUCTURE
# ════════════════════════════════════════════════════════════════════════════
add_rect(slide, C4_X, TOP, C4_W, H,
         RGBColor(0x20, 0x14, 0x02), border=ACCENT_AMBER, border_pt=1.5, corner=8000)
add_rect(slide, C4_X, TOP, C4_W, 0.44,
         ACCENT_AMBER, corner=8000)
add_text(slide, "INFRASTRUCTURE", C4_X+0.08, TOP+0.07, C4_W-0.15, 0.32,
         size=9, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)

def col4_block(y, h, title, lines=None):
    add_rect(slide, C4_X+0.15, y, C4_W-0.30, h,
             RGBColor(0x28, 0x18, 0x04), border=ACCENT_AMBER, border_pt=0.8, corner=6000)
    add_text(slide, title, C4_X+0.24, y+0.05, C4_W-0.50, 0.26,
             size=8.5, bold=True, color=ACCENT_AMBER)
    if lines:
        add_lines(slide, lines, C4_X+0.24, y+0.30, C4_W-0.50, h-0.32,
                  size=7.5, color=LIGHT_GREY)

col4_block(TOP+0.55, 0.80, "NatsConnectionFactory",
           ["Lazy singleton NatsConnection (shared by all services)",
            "BuildOpts()  →  wires CredsFile + auth",
            "IAsyncDisposable  →  disposes on app shutdown"])

col4_block(TOP+1.44, 1.40, "NatsPublisher  (INatsPublisher)",
           ["Serialise  →  Protobuf envelope  +  headers",
            "4-attempt retry:  100 ms → 1 s → 10 s → final",
            "Circuit breaker:  5 failures  →  30 s open window",
            "DLQ route when retries exhausted",
            "Subject: {prefix}.{ctx}.{entity}.{changeType}",
            "Telemetry: ActivitySource  +  Meter counters"])

col4_block(TOP+2.94, 0.64, "NatsCheckpointStore  (ICheckpointStore)",
           ["KV bucket: cm-checkpoints",
            "key: {driverId}.{entityPath}  →  position"])

col4_block(TOP+3.68, 0.75, "FusekiOntologyCache  (IOntologyCache)",
           ["SPARQL query  →  class / property / individual",
            "Indexed by IRI  +  label (case-insensitive)",
            "concept-map  field value  →  ontology IRI"])

col4_block(TOP+4.52, 0.52, "FusekiOntologyCache (HTTP)",
           ["HTTP GET to SPARQL endpoint  •  parses JSON bindings"])

# envelope.proto note
add_rect(slide, C4_X+0.15, TOP+5.16, C4_W-0.30, 0.64,
         RGBColor(0x20, 0x14, 0x02), border=RGBColor(0x80, 0x60, 0x00), border_pt=0.8, corner=6000)
add_text(slide, "envelope.proto  (Protobuf schema)",
         C4_X+0.24, TOP+5.18, C4_W-0.50, 0.26, size=8.5, bold=True, color=ACCENT_AMBER)
add_lines(slide,
    ["EventId (ULID)  •  DriverId  •  Context",
     "EntityPath  •  ChangeType  •  Fields map"],
    C4_X+0.24, TOP+5.44, C4_W-0.50, 0.32, size=7.5)

# ════════════════════════════════════════════════════════════════════════════
# COLUMN 5 — NATS BROKER
# ════════════════════════════════════════════════════════════════════════════
add_rect(slide, C5_X, TOP, C5_W, H,
         RGBColor(0x18, 0x08, 0x00), border=ACCENT_AMBER, border_pt=2.0, corner=8000)
add_rect(slide, C5_X, TOP, C5_W, 0.44,
         ACCENT_AMBER, corner=8000)
add_text(slide, "NATS BROKER", C5_X+0.08, TOP+0.07, C5_W-0.15, 0.32,
         size=9, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)

# JetStream topics
add_rect(slide, C5_X+0.15, TOP+0.55, C5_W-0.30, 0.28,
         RGBColor(0x3C, 0x24, 0x00), corner=4000)
add_text(slide, "JetStream  (persistent)", C5_X+0.20, TOP+0.57, C5_W-0.40, 0.24,
         size=8, bold=True, color=ACCENT_AMBER, align=PP_ALIGN.CENTER)

topics = [
    ("cdc.{ctx}.{entity}.{change}", "Main event stream", ACCENT_AMBER),
    ("cdc.dlq.*",                   "Dead-letter queue",  ACCENT_RED),
    ("cdc.health.{driverId}",       "Health heartbeats",  ACCENT_AMBER),
    ("cdc.lifecycle.{driverId}",    "State transitions",  ACCENT_PURP),
    ("cdc.commands.{driverId}",     "Control commands",   ACCENT_TEAL),
    ("cdc.ontology.refresh",        "Ontology reload",    ACCENT_AMBER),
]
for i, (subj, desc, clr) in enumerate(topics):
    yy = TOP + 0.94 + i * 0.60
    add_rect(slide, C5_X+0.15, yy, C5_W-0.30, 0.54,
             RGBColor(0x24, 0x12, 0x00), border=clr, border_pt=0.8, corner=5000)
    add_text(slide, subj,  C5_X+0.23, yy+0.04, C5_W-0.50, 0.24, size=7.5, bold=True, color=clr)
    add_text(slide, desc,  C5_X+0.23, yy+0.28, C5_W-0.50, 0.20, size=6.5, color=LIGHT_GREY)

# KV store
add_rect(slide, C5_X+0.15, TOP+4.60, C5_W-0.30, 0.28,
         RGBColor(0x3C, 0x24, 0x00), corner=4000)
add_text(slide, "KV Store", C5_X+0.20, TOP+4.62, C5_W-0.40, 0.24,
         size=8, bold=True, color=ACCENT_AMBER, align=PP_ALIGN.CENTER)
add_rect(slide, C5_X+0.15, TOP+5.00, C5_W-0.30, 0.52,
         RGBColor(0x24, 0x12, 0x00), border=ACCENT_AMBER, border_pt=0.8, corner=5000)
add_text(slide, "cm-checkpoints",  C5_X+0.23, TOP+5.02, C5_W-0.50, 0.24, size=8, bold=True, color=ACCENT_AMBER)
add_text(slide, "{driverId}.{entityPath}  →  position",
         C5_X+0.23, TOP+5.24, C5_W-0.50, 0.22, size=7, color=LIGHT_GREY)

# Fuseki
add_rect(slide, C5_X+0.15, TOP+5.65, C5_W-0.30, 0.65,
         RGBColor(0x1A, 0x1A, 0x00), border=YELLOW, border_pt=1.0, corner=5000)
add_text(slide, "Fuseki  (optional external)",
         C5_X+0.23, TOP+5.67, C5_W-0.50, 0.24, size=8, bold=True, color=YELLOW)
add_lines(slide, ["SPARQL endpoint", "Ontology IRI store"],
          C5_X+0.23, TOP+5.90, C5_W-0.50, 0.34, size=7)

# Downstream consumers
add_rect(slide, C5_X+0.15, TOP+6.44, C5_W-0.30, 1.20,
         RGBColor(0x04, 0x1A, 0x04), border=LIME, border_pt=1.2, corner=5000)
add_text(slide, "Downstream Consumers",
         C5_X+0.23, TOP+6.46, C5_W-0.50, 0.26, size=8, bold=True, color=LIME)
add_lines(slide,
    ["Digital twin pipelines",
     "Analytics / dashboards",
     "Event-driven workflows"],
    C5_X+0.23, TOP+6.74, C5_W-0.50, 0.84, size=7.5)

# ════════════════════════════════════════════════════════════════════════════
# FLOW ARROWS between columns
# ════════════════════════════════════════════════════════════════════════════
# Src → Driver
add_arrow(slide, C1_X+C1_W+0.02, TOP+3.50, C2_X-0.02, TOP+3.50,
          color=ACCENT_GREEN, pt=2.0)
# Driver → Host
add_arrow(slide, C2_X+C2_W+0.02, TOP+4.20, C3_X-0.02, TOP+4.20,
          color=ACCENT_TEAL, pt=2.0)
# Host → Infra
add_arrow(slide, C3_X+C3_W+0.02, TOP+3.90, C4_X-0.02, TOP+3.90,
          color=ACCENT_PURP, pt=2.0)
# Infra → NATS
add_arrow(slide, C4_X+C4_W+0.02, TOP+2.90, C5_X-0.02, TOP+2.90,
          color=ACCENT_AMBER, pt=2.0)

# ════════════════════════════════════════════════════════════════════════════
# BOTTOM EVENT-FLOW BANNER
# ════════════════════════════════════════════════════════════════════════════
add_rect(slide, 0.10, 8.60, 15.80, 0.30,
         RGBColor(0x08, 0x08, 0x18), corner=4000)
add_text(slide,
    "Event flow:  Source DB / API  →  IProtocolAdapter  →  RawChangeRecord  →  "
    "FieldMapper  →  RawChangeEvent  →  DefaultEventPipeline  →  "
    "NatsPublisher (protobuf)  →  NATS JetStream  →  Consumers",
    0.18, 8.62, 15.60, 0.28,
    size=7.5, color=ACCENT_TEAL, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════════════════════════════════════
out = r"C:\Repos\UniversalConnector\UniversalConnector_Architecture.pptx"
prs.save(out)
print(f"Saved: {out}")
